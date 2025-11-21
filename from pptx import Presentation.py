from pptx import Presentation
from pptx.util import Pt

slides_data = [
    ('Problem', [
        'Kenyan students face career uncertainty with CBC curriculum.',
        'Lack of personalized guidance for pathway and career selection.',
        'Teachers and parents overwhelmed by CBC complexity.'
    ]),
    ('Solution', [
        'AI-powered platform for CBC career mentorship.',
        'Delivers personalized CBC pathways, career advice, automated CV generation.',
        'Mobile/web access, available 24/7.'
    ]),
    ('Product', [
        'Platform features:',
        '- AI Mentor chatbot',
        '- CBC Pathway Explorer',
        '- AI CV Generator',
        '- Premium: assessment, skill, interview coaching',
        '- Seamless M-Pesa payments'
    ]),
    ('Target Market', [
        'Secondary/high school CBC students',
        'University applicants/job seekers',
        'Parents, guardians, educators',
        'Kenyan youth needing guidance'
    ]),
    ('Market Size', [
        '5M+ CBC students in Kenya',
        '1M university applicants yearly',
        'Total addressable: 7M+ youth and parents'
    ]),
    ('Competitors', [
        'Traditional career offices',
        'Generic mentorship/career apps',
        'Recruitment agencies'
    ]),
    ('Competitive Advantage', [
        'CBC focus: local, contextualized mentorship',
        'AI-driven, always-on guidance',
        'Integrated mobile payments (M-Pesa)',
        'Direct school engagement'
    ]),
    ('Traction & RoadMap', [
        'MVP and dashboard launched',
        'M-Pesa payment pilots',
        'Positive user feedback, strong demand',
        'Next: mobile app, school pilots, 1M user goal'
    ]),
    ('Business/Revenue Model', [
        'Freemium: core features always free',
        'Premium AI tools via subscription (M-Pesa)',
        'School/group licensing',
        'Career package bundles'
    ]),
    ('Go To Market', [
        'School partnerships & teacher engagement',
        'CBC-focused digital social campaigns',
        'NGO/sponsor support',
        'Referral/influencer student challenge'
    ]),
    ('Our Ask', [
        'KES 5M seed investment',
        'Growth, tech, and mobile launch',
        'Content & school onboarding'
    ]),
    ('The Team', [
        'Lead: EdTech/AI Kenya',
        'AI/ML developers',
        'CBC expert teachers',
        'Advisors: school leaders, youth NGOs'
    ]),
]

prs = Presentation()
for (title, bullets) in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for bullet in bullets:
        if not frame.text:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
prs.save('CareerMentor_Pitch_Deck.pptx')